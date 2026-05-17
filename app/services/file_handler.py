# app/services/file_handler.py

import os
import json
import uuid
import textwrap
import time
import requests
import mimetypes
import docx
import google.generativeai as genai
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from googleapiclient.discovery import build

# --- 1. Imports from our modules ---
from ..db.json_manager import find_records, update_record
from ..utils.logger import log_and_print
from ..core.llm_provider import api_key_manager, GEMINI_MODEL_NAME
# Note: GoogleSearch imported lazily in functions to avoid circular imports
# --- Constants & Paths ---
BASE_DIR = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
TEMPLATES_DIR = os.path.join(BASE_DIR, "templates")
OUTPUTS_DIR = os.path.join(BASE_DIR, "outputs")
PERSONAS_DIR = os.path.join(BASE_DIR, "personas")

GOOGLE_API_KEY = "AIzaSyBRhbGOKPT1W3H0GdfaE8X87E2hj1DnOYc" # אותו מפתח של Gemini
SEARCH_ENGINE_ID = "21e7585e5ebf94523"

# Lazy initialization for GoogleSearch to avoid circular imports
_google_search_instance = None

def get_google_search():
    """Get or create GoogleSearch instance lazily."""
    global _google_search_instance
    if _google_search_instance is None:
        from ..services.web_research import GoogleSearch
        _google_search_instance = GoogleSearch(api_key=GOOGLE_API_KEY, search_engine_id=SEARCH_ENGINE_ID)
    return _google_search_instance

# יצירת תיקיות אם לא קיימות
os.makedirs(TEMPLATES_DIR, exist_ok=True)
os.makedirs(OUTPUTS_DIR, exist_ok=True)

# =========================================================================================
#  1. OCR & TEXT EXTRACTION (Gemini Vision + Local Parsers)
# =========================================================================================



def extract_text_from_file(file_path):
    """
    מחלץ תוכן מקבצים.
    עבור טקסט ו-Word: משתמש בחילוץ מקומי.
    עבור PDF ותמונות: משתמש ב-Gemini Vision ומחזיר טקסט מובנה (Markdown).
    """
    import mimetypes
    
    # 1. זיהוי סוג הקובץ
    mime_type, _ = mimetypes.guess_type(file_path)
    _, ext_with_dot = os.path.splitext(file_path)
    ext = ext_with_dot.lower().lstrip('.')
    
    # טיפול מהיר בקבצי טקסט/קוד
    if ext == 'txt':
        try:
            with open(file_path, 'r', encoding='utf-8') as f: return f.read()
        except: pass
            
    # טיפול מקומי ב-Word
    if ext in ['docx', 'doc']:
        doc = docx.Document(file_path)
        return '\n'.join([p.text for p in doc.paragraphs])


    # --- טיפול ב-PDF ותמונות דרך Gemini Vision ---
    from ..core.llm_provider import GEMINI_VISION_MODEL
    
    supported_mime_types = ['application/pdf', 'image/jpeg', 'image/png', 'image/webp', 'image/heic', 'image/heif']

    if not mime_type: # Fallback אם הזיהוי האוטומטי נכשל
        if ext == 'pdf': mime_type = 'application/pdf'
        elif ext in ['jpg', 'jpeg']: mime_type = 'image/jpeg'
        elif ext == 'png': mime_type = 'image/png'
        elif ext == 'webp': mime_type = 'image/webp'

    if mime_type in supported_mime_types:
        log_and_print(f"--- 👁️ Visually analyzing (advanced OCR) the file: '{os.path.basename(file_path)}' ---", "SYSTEM")
        
        try:
            api_key = api_key_manager._get_next_key()
            genai.configure(api_key=api_key)
            
            uploaded_file = genai.upload_file(file_path, mime_type=mime_type)
            
            while uploaded_file.state.name == "PROCESSING":
                time.sleep(1)
                uploaded_file = genai.get_file(uploaded_file.name)
                
            if uploaded_file.state.name == "FAILED":
                raise ValueError("Google failed to process the file.")

            # --- הפרומפט המשופר (The Pro Prompt) ---
            vision_prompt = """
                You are a high-precision OCR and Document Analysis engine. Your task is to convert this document into a structured Markdown representation while strictly preserving the original content and order.

                **CORE INSTRUCTIONS:**
                1. **Linear Processing:** Process the document strictly from start to end. Do not reorder sections.
                2. **Full OCR:** Extract ALL visible text exactly as it appears. Do not summarize or paraphrase.
                3. **Structure:** Use Markdown syntax to represent the document structure:
                - Use headers (#, ##, ###) for titles.
                - Use bullet points (-) or numbered lists (1.) where they appear.
                - Use bold (**) and italics (*) to match emphasis.

                **HANDLING VISUALS (Images, Charts, Diagrams, Handwriting):**
                When you encounter any visual element, insert a description block using exactly this format:

                <figure_description>
                **Type:** [Chart / Diagram / Photo / Handwriting / Form]
                **Visual Content:** Describe exactly what is seen (labels, axes, symbols, layout).
                **Data Extracted:** If data values are shown, list them clearly or use a mini-table.
                **Context:** Briefly state the purpose of this visual as conveyed by the surrounding text.
                </figure_description>

                **HANDLING TABLES:**
                If you encounter a table, DO NOT describe it as an image. instead, reconstruct it as a valid **Markdown Table**. Ensure all rows and columns are preserved.

                **CONSTRAINTS:**
                - If the document is in Hebrew, output strictly in Hebrew (preserve language).
                - Do not output any introductory text like "Here is the analysis". Start directly with the document content.
                - Do not add external knowledge or assumptions.
                """
            
            model = genai.GenerativeModel(GEMINI_VISION_MODEL)
            response = model.generate_content([uploaded_file, vision_prompt])
            extracted_text = response.text

            # ניקוי מהענן
            try: genai.delete_file(uploaded_file.name)
            except: pass

            return extracted_text

        except Exception as e:
            log_and_print(f"Gemini Vision extraction failed: {e}", "ERROR")
            return None

    return None

# =========================================================================================
#  2. MISSION INTELLIGENCE (TIC Updates)
# =========================================================================================

def analyze_document_and_update_tic(agent, mission_id, file_path, filename):
    """
    1. שומר את הקובץ בדיסק (כבר בוצע ב-route, אנחנו מקבלים את הנתיב).
    2. סורק את התוכן המלא ומעדכן את כל ה-TIC: מצב, משימה הבאה, ותזמון.
    3. מוסיף את הקובץ ל-'files_index' ב-TIC לשימוש עתידי.
    """
    # 1. שליפת פרטי המשימה
    missions = find_records(agent, 'missions', {'mission_id': mission_id})
    if not missions: return {"error": "Mission not found"}
    mission = missions[0]
    
    goal = mission.get('goal_description', '')
    current_tic = mission.get('task_information_center', {})
    current_mission_state = current_tic.get('mission_state', {})
    
    # 2. חילוץ טקסט מלא
    file_content = extract_text_from_file(file_path)
    if not file_content:
        return {"error": "Could not extract text from file."}
    
    # אנחנו שולחים את הטקסט לניתוח (עדיין מגבילים למניעת קריסת זיכרון במקרים קיצוניים)
    truncated_content = file_content[:100000] 

    log_and_print(f"--- 📁 File '{filename}' scanned. Performing comprehensive TIC update (State, Action, Schedule)... ---", "SYSTEM")
    
    # 3. הפרומפט המורחב: מעדכן את כל חלקי ה-TIC
    prompt = textwrap.dedent(f"""
        You are the Mission Data Architect & Strategist. A new file has been uploaded.
        
        **Mission Goal:** "{goal}"
        **File Name:** "{filename}"
        
        **Current TIC State:**
        {json.dumps(current_mission_state, ensure_ascii=False)}
        
        **File Content:**
        {truncated_content}

        **YOUR TASKS:**
        Analyze the file content and decide how it changes the ENTIRE mission structure.
        
        1. Update Mission State:
            Extract and consolidate all CRITICAL facts, budget figures, and key entities into mission_state.
            When creating a new mission_state, do not include data from scheduling_preferences or file_description.
            Preserve all existing important information, but restructure and reorganize it to ensure the state is clear, consistent, and logically organized.
            You do not need to include working days in this section.
            
        2. **Update Next Actionable Task:** 
           - Does this file dictate the *immediate* next step? (e.g., "Sign this contract", "Fix these bugs listed in the file").
           - If yes, define the new `next_actionable_task` clearly in Hebrew.
           - If the file is just reference material, you can keep the current task or change it to "לנתח את המידע בקובץ ולקבל החלטות".
           
        3. **Update Scheduling Preferences:**
           - Does the file contain specific working days?
           - If yes, update `scheduling_preferences` (`{{ "work_days": ["Monday", "Wednesday"] }}`).
           - the only param that should be here is "work_days". you dont need to add hours (it should look as example above)
           - If not, return `null` to keep existing preferences.

        4. **Generate File Description:** Write a short 1-sentence summary of the file.

        **OUTPUT FORMAT (JSON ONLY):**
        {{
            "updated_mission_state": {{ ...new state object... }},
            "next_actionable_task": "The new next step (or null to keep current)",
            "scheduling_preferences": {{ ... }} (or null to keep current),
            "file_description": "Short description"
        }}
    """)
    
    try:
        response = api_key_manager.generate_content(prompt, model_name=GEMINI_MODEL_NAME)
        cleaned_response = response.strip().replace("```json", "").replace("```", "")
        result = json.loads(cleaned_response)
        
        # חילוץ הנתונים מהתשובה
        new_state = result.get('updated_mission_state')
        new_next_task = result.get('next_actionable_task')
        new_scheduling = result.get('scheduling_preferences')
        file_desc = result.get('file_description', 'File content')

        # 4. עדכון ה-TIC (רק אם ה-AI החזיר ערכים חדשים)
        if new_state:
            current_tic['mission_state'] = new_state
            
        if new_next_task:
            current_tic['next_actionable_task'] = new_next_task
            log_and_print(f"   - ✅ Next task updated to: {new_next_task}", "SYSTEM")
            
        if new_scheduling:
            current_tic['scheduling_preferences'] = new_scheduling
            log_and_print(f"   - ✅ Scheduling preferences updated based on file.", "SYSTEM")
        
        # עדכון אינדקס הקבצים (files_index)
        if 'files_index' not in current_tic:
            current_tic['files_index'] = []
            
        file_entry = {
            "filename": filename,
            "path": file_path, 
            "description": file_desc,
            "uploaded_at": datetime.now().isoformat()
        }
        
        # הסרה של כפילויות ישנות
        current_tic['files_index'] = [f for f in current_tic['files_index'] if f['filename'] != filename]
        current_tic['files_index'].append(file_entry)

        # שמירה ב-DB
        update_record(agent, 'missions', {'mission_id': mission_id}, {'task_information_center': current_tic})
        
        # הפעלה מחדש של המוח (כדי שיגיב לשינויים החדשים)
        #threading.Thread(target=autonomously_advance_mission, args=(agent, mission_id)).start()
        
        return {
            "status": "Success", 
            "message": "File analyzed. TIC updated (State, Task, Schedule) based on file content.",
            "updated_tic": current_tic # מחזיר את ה-TIC המעודכן לממשק
        }

    except Exception as e:
        log_and_print(f"Error processing file analysis: {e}", "ERROR")
        return {"error": str(e)}


def inspect_file_content(agent, mission_id: str, filename: str, specific_question: str):
    """
    כלי אחזור (Retrieval): קורא את התוכן המלא של קובץ שמור מהדיסק,
    מנתח אותו מול שאלה ספציפית, ומחזיר רק את התשובה הרלוונטית.
    """
    log_and_print(f"--- 🔍 Opening file '{filename}' to answer: '{specific_question}' ---", "SYSTEM")
    
    # 1. שליפת מיקום הקובץ מתוך ה-TIC
    missions = find_records(agent, 'missions', {'mission_id': mission_id})
    if not missions: return {"error": "Mission not found"}
    
    tic = missions[0].get('task_information_center', {})
    files_index = tic.get('files_index', [])
    
    # מציאת הקובץ באינדקס
    file_entry = next((f for f in files_index if f['filename'] == filename), None)
    
    if not file_entry:
        # נסיון חיפוש גמיש (אם ה-AI לא דייק בשם הקובץ)
        file_entry = next((f for f in files_index if filename.lower() in f['filename'].lower()), None)
        
    if not file_entry:
        return {"error": f"File '{filename}' not found in mission index."}
    
    file_path = file_entry.get('path')
    if not os.path.exists(file_path):
        return {"error": "File record exists, but physical file is missing from server."}
        
    # 2. קריאת התוכן המלא בזמן אמת
    full_text = extract_text_from_file(file_path)
    if not full_text:
        return {"error": "File is empty or unreadable."}
        
    # 3. שליחה ל-AI (Q&A)
    # אנו שולחים רק את השאלה והתוכן. זה לא נכנס ללוג הראשי של המשימה!
    qa_prompt = textwrap.dedent(f"""
        You are an intelligent file analyzer.
        
        **User Question:** "{specific_question}"
        
        **File Content:**
        {full_text[:100000]}
        
        **Instructions:**
        1. Answer the specific question based ONLY on the file content.
        2. If the answer is not in the file, state that clearly.
        3. Keep the answer concise and direct.
        4. Answer in Hebrew.
    """)
    
    try:
        response = api_key_manager.generate_content(qa_prompt, model_name=GEMINI_MODEL_NAME)
        answer = response.strip()
        
        log_and_print(f"   - ✅ Answer extracted from file: {answer[:50]}...", "SYSTEM")
        
        # מחזירים רק את התשובה. זה מה שייכנס להיסטוריה הקצרה, אבל לא יעמיס על המערכת
        return {"status": "Success", "answer_from_file": answer}
        
    except Exception as e:
        return {"error": f"Error analyzing file content: {e}"}

# =========================================================================================
#  3. POWERPOINT GENERATION
# =========================================================================================

def generate_presentation_feature(agent, topic: str, content_guidelines: str, template_name: str = "template1.pptx"):
    """
    גרסה "ראש גדול": אם אין מקום מוגדר לטקסט, הבוט יוצר אותו בעצמו.
    מותאם במיוחד לרקע כהה ועברית.
    """
    # הגנה מסוגי משתנים
    if isinstance(template_name, list): template_name = template_name[0]
    
    log_and_print(f"--- 🎨 Building presentation (free mode) on topic: '{topic}' ---", "SYSTEM")
    
    # הגדרת נתיבים
    templates_dir = os.path.join(BASE_DIR, "templates")
    output_dir = os.path.join(BASE_DIR, "outputs")
    os.makedirs(output_dir, exist_ok=True)
    
    template_path = os.path.join(templates_dir, template_name)
    if not os.path.exists(template_path):
        # נסיון למצוא קובץ כלשהו
        files = [f for f in os.listdir(templates_dir) if f.endswith('.pptx')]
        if files: template_path = os.path.join(templates_dir, files[0])
        else: return {"error": "No PPTX file found in templates folder."}

    try:
        # שלב 1: תכנון המצגת ע"י ה-AI
        planning_prompt = textwrap.dedent(f"""
            Create a structure for a PowerPoint presentation.
            Topic: "{topic}"
            Guidelines: "{content_guidelines}"
            
            Output JSON format:
            {{
                "slides": [
                    {{
                        "type": "title_slide",
                        "title": "Main Title Here",
                        "subtitle": "Subtitle Here"
                    }},
                    {{
                        "type": "content_slide",
                        "title": "Slide Title",
                        "points": ["Bullet point 1", "Bullet point 2", "Bullet point 3"],
                        "image_query": "Optional image search term"
                    }}
                ]
            }}
        """)

        response = api_key_manager.generate_content(planning_prompt, model_name=GEMINI_MODEL_NAME)
        
        # ניקוי התשובה
        txt = response.strip()
        if "```json" in txt: txt = txt.split("```json")[1].split("```")[0]
        elif "```" in txt: txt = txt.split("```")[1].split("```")[0]
        plan = json.loads(txt.strip())

        # שלב 2: בניית המצגת
        prs = Presentation(template_path)
        
        # מחיקת שקפים קיימים כדי להתחיל נקי
        xml_slides = prs.slides._sldIdLst  
        slides_list = list(xml_slides)
        for s in slides_list: xml_slides.remove(s)

        # נשתמש ב-Layout הראשון (בדרך כלל שומר על הרקע)
        master_layout = prs.slide_layouts[0] 

        temp_images = []

        for slide_data in plan.get("slides", []):
            slide = prs.slides.add_slide(master_layout)
            
            # --- כותרת ---
            title_text = slide_data.get("title", "")
            
            # אם יש כותרת מובנית בתבנית - נשתמש בה
            if slide.shapes.title:
                slide.shapes.title.text = title_text
            else:
                # אם אין, ניצור תיבת כותרת למעלה
                left = Inches(0.5)
                top = Inches(0.5)
                width = Inches(9)
                height = Inches(1.5)
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                tf.text = title_text
                p = tf.paragraphs[0]
                p.font.size = Pt(44)
                p.font.bold = True
                p.font.name = "Arial"
                p.font.color.rgb = RGBColor(255, 255, 255) # לבן
                p.alignment = PP_ALIGN.CENTER

            # --- תוכן (Subtitle או נקודות) ---
            body_text = ""
            if "subtitle" in slide_data:
                body_text = slide_data["subtitle"]
            elif "points" in slide_data:
                body_text = "\n".join([f"• {p}" for p in slide_data["points"]])

            if body_text:
                # יצירת תיבת טקסט ידנית לגוף המצגת
                left = Inches(1)
                top = Inches(2.5)
                width = Inches(8)
                height = Inches(4)
                
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                tf.word_wrap = True
                
                # הכנסת הטקסט
                tf.text = body_text
                
                # עיצוב הטקסט (חשוב!)
                for p in tf.paragraphs:
                    p.font.size = Pt(24)
                    p.font.name = "Arial"
                    p.font.color.rgb = RGBColor(255, 255, 255) # לבן לרקע שחור
                    p.alignment = PP_ALIGN.RIGHT # יישור לימין לעברית
                    p.space_after = Pt(14) # רווח בין שורות

            # --- תמונה ---
            if "image_query" in slide_data:
                q = slide_data["image_query"]
                if isinstance(q, list): q = " ".join(q)
                img_path = get_image_for_slide(str(q))
                if img_path:
                    temp_images.append(img_path)
                    # ממקם את התמונה בצד שמאל למטה
                    slide.shapes.add_picture(img_path, Inches(0.5), Inches(4), height=Inches(3))

        # שמירה
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        safe_topic = "".join([c for c in topic if c.isalpha() or c.isdigit() or c==' ']).strip().replace(" ", "_")
        filename = f"Presentation_{safe_topic}_{timestamp}.pptx"
        output_path = os.path.join(output_dir, filename)
        prs.save(output_path)
        
        # ניקוי
        for img in temp_images:
            try: os.remove(img)
            except: pass
            
        return {
            "status": "Success",
            "message": f"Presentation generated successfully: {filename}",
            "download_link": f"/outputs/{filename}"
        }

    except Exception as e:
        log_and_print(f"ERROR: {e}", "ERROR")
        return {"error": str(e)}


def get_image_for_slide(query):
    """מחפש תמונה, מוריד אותה זמנית ומחזיר נתיב. חסין לשגיאות טייפ."""
    try:
        # הגנה: אם התקבל list במקום string, נמיר אותו
        if isinstance(query, list):
            query = " ".join(query)
        if not isinstance(query, str):
            return None

        log_and_print(f"   🖼️ Searching for image: {query}", "SYSTEM")

        # חיפוש
        search_results = get_google_search().search([f"{query}"], num_results=2)
        image_url = None
        
        for res_set in search_results:
            for item in res_set.get('results', []):
                link = item.get('link', '')
                # בדיקה בסיסית לסיומת תמונה
                if link.lower().endswith(('.jpg', '.jpeg', '.png', '.webp')):
                    image_url = link
                    break
            if image_url: break
        
        # Fallback ל-Unsplash אם לא נמצאה תמונה
        if not image_url:
            safe_query = query.replace(' ', '-')
            image_url = f"https://source.unsplash.com/featured/?{safe_query}"

        # הורדה
        response = requests.get(image_url, stream=True, timeout=5)
        if response.status_code == 200:
            temp_filename = f"temp_img_{uuid.uuid4().hex}.jpg"
            temp_path = os.path.join(BASE_DIR, temp_filename)
            with open(temp_path, 'wb') as out_file:
                out_file.write(response.content)
            return temp_path
            
    except Exception as e:
        print(f"Image download warning: {e}")
        return None



